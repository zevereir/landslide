from pptx import Presentation
from pptx.shapes.placeholder import PicturePlaceholder, BasePlaceholder, TablePlaceholder
from PIL import Image
from thesis_sieben_bocklandt.code.prototyping.classes import *
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_AUTO_SIZE

# Enumerates voor de verschillende slide layouts
TITLE_SLIDE = 0
SINGLE_CONTENT = 1
SECTION_HEADER = 2
DOUBLE_CONTENT = 3
TRIPLE_CONTENT = 4
COMPARISON = 5
TITLE_ONLY = 6
CAPTIONED_CONTENT = 7
BACKGROUND_QUOTE = 8
BACKGROUND_ONLY = 9
CONTENT_ONLY = 10
IMAGES = ["picture", "combined_pictures", "background"]


def archetypes2slides(archetypes, tree, output_directory, ppt_path, RA=None, make_slides=True):
    """
    De functie die gebruik makend van de archetypes en de categorized xml_tree een nieuwe layout-loze
    powerpoint opmaakt. De manier waarop een slide wordt opgebouwd hangt af van het archetype.
    """
    prs = Presentation(ppt_path)
    footer_dimensions = []
    slide_number_dimensions = []
    placeholder_percentage = []
    for l in prs.slide_layouts:
        for pl in l.placeholders:
            if pl.placeholder_format.idx == 10 and footer_dimensions == []:
                footer_font = Pt(10)
                footer_dimensions = [pl.top, pl.left, pl.width, pl.height]
            elif pl.placeholder_format.idx == 11 and slide_number_dimensions == []:
                slide_number_dimensions = [pl.top, pl.left, pl.width, pl.height]
                slide_number_font = Pt(10)
    info_used = []
    all_archetypes = archetypes[:]
    references = []
    slide_numbers = []
    for page in tree.getroot():
        nb_elements = len(page)
        slide_numbers.append(len(page.findall("slide_number")) != 0)
        used_content = []
        good_elements = [x for x in page if x.attrib.get("archetypal_index") != None]
        archetype = all_archetypes.pop(0)
        page_size = [float(x) for x in page.attrib.get("bbox").split(",")]
        slide_height = prs.slide_height
        slide_width = prs.slide_width
        to_remove_as_good_element = set()
        if isinstance(archetype, TitleSlide):
            title = good_elements[archetype.title].text
            to_remove_as_good_element.add(archetype.title)
            subscript = []
            for i in archetype.subscript:
                subscript.append(good_elements[i])
                if good_elements[i].tag not in IMAGES:
                    to_remove_as_good_element.add(i)
            placeholder_percentage.append((1 + len(subscript)) / nb_elements)
            slide = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE])
            slide_title_slide(title, subscript, slide, used_content)
        elif isinstance(archetype, TitleSingleContent):
            placeholder_percentage.append(2 / nb_elements)
            title = good_elements[archetype.title].text
            content = good_elements[archetype.content]
            to_remove_as_good_element.add(archetype.title)
            to_remove_as_good_element.add(archetype.content)
            slide = prs.slides.add_slide(prs.slide_layouts[SINGLE_CONTENT])
            slide_single_content(title, content, page_size, slide_height, slide_width, slide, output_directory,
                                 used_content)
        elif isinstance(archetype, SectionHeader):
            title = good_elements[archetype.title].text
            to_remove_as_good_element.add(archetype.title)
            subscript = []
            for i in archetype.subscript:
                subscript.append(good_elements[i])
                if good_elements[i].tag not in IMAGES:
                    to_remove_as_good_element.add(i)
            placeholder_percentage.append((1 + len(subscript)) / nb_elements)
            slide = prs.slides.add_slide(prs.slide_layouts[SECTION_HEADER])
            slide_section_header(title, subscript, slide, used_content)
        elif isinstance(archetype, TitleDoubleContent):
            placeholder_percentage.append(3 / nb_elements)
            title = good_elements[archetype.title].text
            to_remove_as_good_element.add(archetype.title)
            to_remove_as_good_element.add(archetype.left_content)
            to_remove_as_good_element.add(archetype.right_content)

            left_content = good_elements[archetype.left_content]
            right_content = good_elements[archetype.right_content]
            slide = prs.slides.add_slide(prs.slide_layouts[DOUBLE_CONTENT])
            slide_double_content(title, left_content, right_content, page_size, slide_height, slide_width, slide,
                                 output_directory, used_content)
        elif isinstance(archetype, TitleTripleContent):
            placeholder_percentage.append(4 / nb_elements)
            title = good_elements[archetype.title].text
            left_content = good_elements[archetype.left_content]
            middle_content = good_elements[archetype.middle_content]
            right_content = good_elements[archetype.right_content]
            to_remove_as_good_element.add(archetype.title)
            to_remove_as_good_element.add(archetype.left_content)
            to_remove_as_good_element.add(archetype.middle_content)
            to_remove_as_good_element.add(archetype.right_content)
            slide = prs.slides.add_slide(prs.slide_layouts[TRIPLE_CONTENT])
            slide_triple_content(title, left_content, middle_content, right_content, page_size, slide_height,
                                 slide_width,
                                 slide, output_directory, used_content)
        elif isinstance(archetype, Comparison):
            placeholder_percentage.append(5 / nb_elements)
            title = good_elements[archetype.title].text
            left_subtitle = good_elements[archetype.left_subtitle]
            right_subtitle = good_elements[archetype.right_subtitle]
            left_content = good_elements[archetype.left_content]
            right_content = good_elements[archetype.right_content]
            to_remove_as_good_element.add(archetype.title)
            to_remove_as_good_element.add(archetype.left_subtitle)
            to_remove_as_good_element.add(archetype.right_subtitle)
            to_remove_as_good_element.add(archetype.left_content)
            to_remove_as_good_element.add(archetype.right_content)
            slide = prs.slides.add_slide(prs.slide_layouts[COMPARISON])
            slide_comparison(title, left_subtitle, right_subtitle, left_content, right_content, page_size, slide_height,
                             slide_width, slide, output_directory, used_content)
        elif isinstance(archetype, TitleOnly):
            placeholder_percentage.append(1 / nb_elements)
            title = good_elements[archetype.title].text
            to_remove_as_good_element.add(archetype.title)
            slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY])
            slide_title_only(title, slide, used_content)
        elif isinstance(archetype, CaptionedContent):
            placeholder_percentage.append(3 / nb_elements)
            top_content = good_elements[archetype.top_content]
            left_content = good_elements[archetype.left_content]
            right_content = good_elements[archetype.right_content]
            to_remove_as_good_element.add(archetype.top_content)
            to_remove_as_good_element.add(archetype.left_content)
            to_remove_as_good_element.add(archetype.right_content)
            slide = prs.slides.add_slide(prs.slide_layouts[CAPTIONED_CONTENT])
            slide_captioned_content(top_content, left_content, right_content, slide, output_directory, slide_height,
                                    slide_width, page_size, used_content)
        elif isinstance(archetype, BackgroundQuote):
            placeholder_percentage.append(2 / nb_elements)
            title = good_elements[archetype.title].text
            to_remove_as_good_element.add(archetype.title)
            background = good_elements[archetype.background]
            to_remove_as_good_element.add(archetype.background)
            slide = prs.slides.add_slide(prs.slide_layouts[BACKGROUND_QUOTE])
            slide_background_quote(title, background, slide, output_directory, slide_height, slide_width, page_size,
                                   used_content)
        elif isinstance(archetype, BackgroundOnly):
            placeholder_percentage.append(1 / nb_elements)
            background = good_elements[archetype.background]
            to_remove_as_good_element.add(archetype.background)
            slide = prs.slides.add_slide(prs.slide_layouts[BACKGROUND_ONLY])
            slide_background_only(background, slide, output_directory, slide_height, slide_width, page_size,
                                  used_content)
        elif isinstance(archetype, ContentOnly):
            placeholder_percentage.append(0 / max(1, nb_elements))
            contents = [good_elements[x] for x in archetype.content]
            for element in archetype.content:
                to_remove_as_good_element.add(element)
            slide = prs.slides.add_slide(prs.slide_layouts[CONTENT_ONLY])
            slide_content_only(contents, prs.slide_width, prs.slide_height, page_size, slide, output_directory,
                               used_content)
        else:
            if len(archetype) > 1:
                if isinstance(archetype[0], ContentOnly):
                    placeholder_percentage.append(0 / max(1, nb_elements))
                    contents = [good_elements[x] for x in archetype.content]
                    for element in archetype.content:
                        to_remove_as_good_element.add(element)
                    slide = prs.slides.add_slide(prs.slide_layouts[CONTENT_ONLY])
                    slide_content_only(contents, prs.slide_width, prs.slide_height, page_size, slide, output_directory,
                                       used_content)
            raise TypeError
        # if RA!=None:
        #     if len(good_elements) != len(to_remove_as_good_element):
        #         print("Written archetype2slides")
        #         with open(
        #                 "D:\\User\\Documents\\School\\Thesis\\thesis_sieben_bocklandt\\Data\\clustering\\clusters.txt",'a') as f:
        #             if RA[int(page.attrib.get("id"))-1][0] !=0:
        #                 f.write(str(RA[int(page.attrib.get("id"))-1])+"\n")
        for i in range(0, len(good_elements)):
            if i not in to_remove_as_good_element:
                content = good_elements[i]
                bbox = [float(x) for x in content.attrib.get("bbox").split(",")]
                x_factor = slide_width / page_size[2]
                y_factor = slide_height / page_size[3]
                add_content_without_placeholder(content, bbox, slide, output_directory, slide_height, x_factor,
                                                y_factor, used_content)

        reference = " "
        for ref in page.findall("ref"):
            reference += (reference != " ") * "\n" + ref.text
            used_content.append(ref.text)
        references.append(reference)
        to_remove = set()
        to_add = set()

        for i in used_content:
            if i != None and "\n" in i:
                to_remove.add(i)
                for x in i.split("\n"):
                    to_add.add(x)
        for i in to_remove:
            used_content = [x for x in used_content if x != i]
        for i in to_add:
            used_content.append(i)
        info_used.append([x.strip() for x in used_content if x != None])
    add_numbers = sum(slide_numbers) / len(slide_numbers) > 0.8
    counter = 0
    for slide in prs.slides:
        ref = references.pop(0)
        txBox = slide.shapes.add_textbox(footer_dimensions[1], footer_dimensions[0], footer_dimensions[2],
                                         footer_dimensions[3])
        tf = txBox.text_frame
        tf.text = ref
        tf.paragraphs[0].font.size = footer_font
        if add_numbers:
            if (counter == 0 and slide_numbers[0]) or counter > 0:
                txBox_slidenumber = slide.shapes.add_textbox(slide_number_dimensions[1], slide_number_dimensions[0],
                                                             slide_number_dimensions[2],
                                                             slide_number_dimensions[3])
                tf_slidenumber = txBox_slidenumber.text_frame
                tf_slidenumber.text = str(counter + slide_numbers[0])
                tf_slidenumber.paragraphs[0].font.size = slide_number_font
                info_used[counter] = info_used[counter] + [str(counter + slide_numbers[0])]
            counter += 1
    if make_slides:
        prs.save(output_directory + "\\test.pptx")
    info_used=[x for x in info_used if x!=None]
    return list(zip(info_used, placeholder_percentage))


def slide_title_slide(titletext, content, slide, used_content):
    """
    De functie die een title slide opbouwt
    """
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text_frame.word_wrap = True
    #title.text_frame.text = titletext
    used_content.append(titletext)
    subtext = ""
    for i in content:
        if i.tag not in IMAGES:
            if i.text != None:
                subtext += i.text + "\n"
            used_content.append(i.text)
    tf = subtitle.text_frame
    tf.text = subtext


def slide_section_header(titletext, content, slide, used_content):
    """
        De functie die een section header slide opbouwt
    """
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    #title.text = titletext
    used_content.append(titletext)
    subtext = ""
    for i in content:
        if i.tag not in IMAGES:
            if i.text != None:
                subtext += i.text + "\n"
                used_content.append(i.text)
    subtitle.text = subtext[:-1]


def slide_single_content(titletext, content, page_size, slide_height, slide_width, slide, output_directory,
                         used_content):
    """
        De functie die een single content slide opbouwt
    """
    title = slide.shapes.title
    single_content = slide.placeholders[1]
   # title.text = titletext
    used_content.append(titletext)
    add_content(slide, single_content, content, output_directory, slide_height, slide_width, page_size, used_content)


def slide_double_content(titletext, left_content, right_content, page_size, slide_height, slide_width, slide,
                         output_directory, used_content):
    """
            De functie die een dubbele content slide opbouwt
    """
    title = slide.shapes.title
    left_content_placeholder = slide.placeholders[1]
    right_content_placeholder = slide.placeholders[2]
   # title.text = titletext
    used_content.append(titletext)
    add_content(slide, left_content_placeholder, left_content, output_directory, slide_height, slide_width, page_size,
                used_content)
    add_content(slide, right_content_placeholder, right_content, output_directory, slide_height, slide_width, page_size,
                used_content)


def slide_triple_content(titletext, left_content, middle_content, right_content, page_size, slide_height, slide_width,
                         slide, output_directory, used_content):
    """
    De functie die een triple content slide opbouwt
    """

    title = slide.shapes.title
    #title.text = titletext
    used_content.append(titletext)
    left_content_placeholder = slide.placeholders[1]
    middle_content_placeholder = slide.placeholders[13]
    right_content_placeholder = slide.placeholders[14]
    add_content(slide, left_content_placeholder, left_content, output_directory, slide_height, slide_width, page_size,
                used_content)
    add_content(slide, middle_content_placeholder, middle_content, output_directory, slide_height, slide_width,
                page_size, used_content)
    add_content(slide, right_content_placeholder, right_content, output_directory, slide_height, slide_width, page_size,
                used_content)


def slide_comparison(titletext, left_subtitle, right_subtitle, left_content, right_content, page_size, slide_height,
                     slide_width, slide, output_directory, used_content):
    """
        De functie die een Comparison slide opbouwt
    """
    title = slide.shapes.title
    #title.text = titletext
    used_content.append(titletext)
    left_subtitle_placeholder = slide.placeholders[2]
    right_subtitle_placeholder = slide.placeholders[4]
    left_content_placeholder = slide.placeholders[1]
    right_content_placeholder = slide.placeholders[3]
    add_content(slide, left_subtitle_placeholder, left_subtitle, output_directory, slide_height, slide_width, page_size,
                used_content)
    add_content(slide, right_subtitle_placeholder, right_subtitle, output_directory, slide_height, slide_width,
                page_size, used_content)
    add_content(slide, left_content_placeholder, left_content, output_directory, slide_height, slide_width, page_size,
                used_content)
    add_content(slide, right_content_placeholder, right_content, output_directory, slide_height, slide_width, page_size,
                used_content)


def slide_title_only(titletext, slide, used_content):
    """
    De functie die in title-only slide opbouwt
    """
    title = slide.shapes.title
   # title.text = titletext
    used_content.append(titletext)


def slide_captioned_content(top_content, left_content, right_content, slide, output_directory, slide_height,
                            slide_width, page_size, used_content):
    """
    De functie die een captioned content slide opbouwt
    """
    top_content_placeholder = slide.placeholders[0]
    left_content_placeholder = slide.placeholders[2]
    right_content_placeholder = slide.placeholders[1]
    add_content(slide, top_content_placeholder, top_content, output_directory, slide_height, slide_width, page_size,
                used_content)
    add_content(slide, left_content_placeholder, left_content, output_directory, slide_height, slide_width, page_size,
                used_content)
    add_content(slide, right_content_placeholder, right_content, output_directory, slide_height, slide_width, page_size,
                used_content)


def slide_background_quote(titletext, background, slide, output_directory, slide_height, slide_width, page_size,
                           used_content):
    """
    De functie die een background-quote opbouwt
    """
    title = slide.shapes.title
    #title.text = titletext
    used_content.append(titletext)
    background_placeholder = slide.placeholders[13]
    add_content(slide, background_placeholder, background, output_directory, slide_height, slide_width, page_size,
                used_content)


def slide_background_only(background, slide, output_directory, slide_height, slide_width, page_size, used_content):
    """
    De functie die een background-only slide opbouwt
    """
    background_placeholder = slide.placeholders[13]
    add_content(slide, background_placeholder, background, output_directory, slide_height, slide_width, page_size,
                used_content)


def slide_content_only(contents, slide_width, slide_height, page_size, slide, output_directory, used_content):
    """
        De functie die een content-only slide opbouwt
    """
    x_factor = slide_width / page_size[2]
    y_factor = slide_height / page_size[3]
    for content in contents:
        add_content_without_placeholder(content, [float(x) for x in content.attrib.get("bbox").split(",")], slide,
                                        output_directory, slide_height, x_factor, y_factor, used_content)


def add_content_without_placeholder(content, bbox, slide, output_directory, slide_height, x_factor, y_factor,
                                    used_content):
    """
    De functie die in een slide een tekstvak toevoegt gebaseerd op de relatieve ligging van de content in de oude slide.
    Deze coordinaten worden geschaald naar de grootte van de slide en op die plaats wordt het tekstvak toegevoegd.
    """

    if content.tag == "picture" or content.tag == "background":

        # slide.shapes.add_picture(output_directory + "\\images\\" + content.attrib.get("src"), bbox[0] * x_factor,
        #                          slide_height-bbox[3] * y_factor, int(x_factor * (bbox[2] - bbox[0])), int(y_factor * (bbox[3] - bbox[1])))
        combined = content.attrib.get("combined_sources")
        if combined == None:
            used_content.append(content.attrib.get("src"))
        else:
            for i in combined[1:-1].split(","):
                used_content.append(i.strip()[1:-1])
        caption = content.attrib.get("caption")
        if caption != None and caption != "":
            # txBox = slide.shapes.add_textbox(x_factor * bbox[0], slide_height - y_factor * bbox[3]+y_factor * (bbox[3] - bbox[1]),
            #                                  x_factor * (bbox[2] - bbox[0]), 0.1*y_factor * (bbox[3] - bbox[1]))
            # tf = txBox.text_frame
            # tf.text = caption
            # for para in tf.paragraphs:
            #     para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            used_content.append(caption)

    elif content.tag == "normal_text" or content.tag == "title":
        # txBox = slide.shapes.add_textbox(x_factor * bbox[0], slide_height - y_factor * bbox[3],
        #                                  x_factor * (bbox[2] - bbox[0]), y_factor * (bbox[3] - bbox[1]))
        # tf = txBox.text_frame
        # tf.text = content.text
        used_content.append(content.text)
    elif content.tag == "table":
        used_content.append("table_bbox_for_score_calc_" + str(bbox))
        dimensions = content.attrib.get("dimensions").split("/")
        table_width = int(dimensions[0])
        table_height = int(dimensions[1])
        table = slide.shapes.add_table(table_height, table_width, int(bbox[0] * x_factor),
                                       int(slide_height - y_factor * bbox[3]),
                                       int(x_factor * (bbox[2] - bbox[0])), int(y_factor * (bbox[3] - bbox[1]))).table
        counter = 0
        for cell in content:

            tablecell = table.cell(counter // table_width, counter % table_width)
            if len(cell) > 0:
                tabletext = cell[0].text
                tablecell.text = tabletext
                used_content.append(tabletext)
            counter += 1

    elif content.tag == "enlisting":
        levels = []
        for line in content:
            line_bbox = [float(x) for x in line.attrib.get("bbox").split(",")]
            if line_bbox[0] not in levels:
                levels.append(line_bbox[0])
        levels.sort()
        txBox = slide.shapes.add_textbox(x_factor * bbox[0], slide_height - y_factor * bbox[3],
                                         x_factor * (bbox[2] - bbox[0]), y_factor * (bbox[3] - bbox[1]))
        tf = txBox.text_frame
        content_list = content[:]
        content_list.sort(reverse=True, key=get_first_y)
        for line in content_list:
            p = tf.add_paragraph()
            line_bbox = [float(x) for x in line.attrib.get("bbox").split(",")]
            level = 1 + levels.index(line_bbox[0])
            enlisting_text = line.text[0 + (ord(line.text[0]) > 128):]
            p.text = enlisting_text
            p.level = min(level, 8)
            used_content.append(line.text)
    else:
        raise TypeError(content.tag)


def add_content(slide, placeholder, content, output_directory, slide_height, slide_width, page_size, used_content):
    """
    De functie die content toevoegt aan een al bestaande placeholder. Deze placeholder bepaald zelf waar de content
    wordt gezet. Het voordeel hiervan is dat deze dan mee gaat verschuiven als er een andere slide master gebruikt wordt
    """
    if content.tag == "picture" or content.tag == "background":
        # source=output_directory+"\\images\\"+content.attrib.get("src")
        combined = content.attrib.get("combined_sources")
        if combined == None:
            used_content.append(content.attrib.get("src"))
        else:
            for i in combined[1:-1].split(","):
                used_content.append(i.strip()[1:-1])
        caption = content.attrib.get("caption")
        if caption != None:
            used_content.append(caption)
        # add_image(slide,placeholder,source, content.tag=="background", content.attrib.get("caption"))
    else:
        if content.tag == "normal_text" or content.tag == "title":
            begin_text = content.text
            used_content.append(content.text)
            # add_text(placeholder,begin_text)
        elif content.tag == "table":
            add_table(slide, placeholder, content, used_content)

        elif content.tag == "enlisting":
            add_enlisting(placeholder, content, used_content)
        else:
            print(content.tag)
            raise TypeError(content.tag)


def add_text(placeholder, text):
    """
    functie om tekst toe te voegen aan een placeholder
    """
    tf = placeholder.text_frame
    tf.text = text
    for para in tf.paragraphs:
        para.font.size = Pt(20)


def add_enlisting(placeholder, content, used_content):
    """"
    functie om een opsomming toe te voegen aan een  placeholder
    """
    levels = []
    for line in content:
        line_bbox = [float(x) for x in line.attrib.get("bbox").split(",")]
        if line_bbox[0] not in levels:
            levels.append(line_bbox[0])
    levels.sort()
    tf = placeholder.text_frame
    content_list = content[:]
    content_list.sort(reverse=True, key=get_first_y)
    for line in content_list:
        p = tf.add_paragraph()
        line_bbox = [float(x) for x in line.attrib.get("bbox").split(",")]
        level = 1 + levels.index(line_bbox[0])
        enlisting_text = line.text[0 + (ord(line.text[0]) > 128):]
        p.text = enlisting_text
        used_content.append(enlisting_text)
        p.level = min(level, 8)


def get_first_y(line):
    """"
    helpfunctie om de hoogte van een tekstlijn te bepalen
    """
    return [float(x) for x in line.attrib.get("bbox").split(",")][1]


def add_image(slide, placeholder, source, background, caption):
    """"
    functie om een afbeelding toe te voegen aan een al bestaande placeholder
    """
    single_content = PicturePlaceholder(placeholder._sp, slide)
    width, height = Image.open(source).size
    # Make sure the placeholder doesn't zoom in
    single_content.height = height
    single_content.width = width
    single_content = single_content.insert_picture(source)
    # Calculate ratios and compare
    if not background:
        image_ratio = width / height
        placeholder_ratio = single_content.width / single_content.height
        ratio_difference = placeholder_ratio - image_ratio
        # Placeholder width too wide:
        if ratio_difference > 0:
            difference_on_each_side = ratio_difference / 2
            single_content.crop_left = -difference_on_each_side
            single_content.crop_right = -difference_on_each_side
        # Placeholder height too high
        else:
            difference_on_each_side = -ratio_difference / 2
            single_content.crop_bottom = -difference_on_each_side
            single_content.crop_top = -difference_on_each_side
    width = single_content.width
    height = single_content.height
    left = single_content.left
    top = single_content.top + height
    if caption != None and caption != "":
        txBox = slide.shapes.add_textbox(left, top, width, 0.1 * height)
        tf = txBox.text_frame
        tf.text = caption
        for para in tf.paragraphs:
            para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER


def add_table(slide, placeholder, content, used_content):
    """"
    functie om een tabel toe te voegen aan een al bestaande placeholder
    """
    dimensions = content.attrib.get("dimensions").split("/")
    used_content.append("table_bbox_for_score_calc_[" + str(content.attrib.get("bbox")) + "]")
    table_width = int(dimensions[0])
    table_height = int(dimensions[1])
    table_placeholder = TablePlaceholder(placeholder._sp, slide)
    graphic_frame = table_placeholder.insert_table(table_height, table_width)
    table = graphic_frame.table
    counter = 0
    for cell in content:
        tablecell = table.cell(counter // table_width, counter % table_width)
        if len(cell) > 0:
            tabletext = cell[0].text
            tablecell.text = tabletext
            used_content.append(tabletext)
        counter += 1

